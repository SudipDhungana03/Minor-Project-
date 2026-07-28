from rest_framework.response import Response
from rest_framework.decorators import api_view
from .detector_service import run_analysis, run_source_verification
from ..models import DetectionResult
from apps.classroom.models import Submission
from .plagiarism_vector import build_similarity_report, extract_text_from_submission_file
import os
import logging
import io
import requests
from threading import Thread

logger = logging.getLogger(__name__)


def _extract_text_from_file(file_field):
    """Try to extract text from common document types. Returns string or None on failure.

    Fallbacks:
    - If the local `file_field.path` does not exist but `file_field.url` is an HTTP(S) URL,
      attempt to download the file into memory and extract from the bytes.
    """
    if not file_field:
        return None

    storage = getattr(file_field, 'storage', None)
    name = getattr(file_field, 'name', None)
    path = getattr(file_field, 'path', None)
    url = getattr(file_field, 'url', None)

    # If local file exists, use it. Otherwise, attempt storage or URL fallback.
    use_bytes = None
    source = None
    if path and os.path.exists(path):
        source = path
    elif storage and name:
        try:
            if storage.exists(name):
                with storage.open(name, 'rb') as f:
                    use_bytes = io.BytesIO(f.read())
        except Exception:
            logger.exception('Failed to open file from storage: %s', name)

    if source is None and use_bytes is None and url and (url.startswith('http://') or url.startswith('https://')):
        try:
            resp = requests.get(url, timeout=10)
            if resp.status_code == 200:
                use_bytes = io.BytesIO(resp.content)
            else:
                logger.warning('Failed to download file for extraction: %s (status %s)', url, resp.status_code)
        except Exception:
            logger.exception('HTTP download for extraction failed for %s', url)

    # determine extension
    ext = None
    if source:
        ext = os.path.splitext(source)[1].lower()
    elif name:
        ext = os.path.splitext(name.split('?')[0])[1].lower()
    elif url:
        ext = os.path.splitext(url.split('?')[0])[1].lower()

    try:
        # PDF extraction
        if ext == '.pdf':
            try:
                from pypdf import PdfReader
                if use_bytes:
                    reader = PdfReader(use_bytes)
                else:
                    reader = PdfReader(source)
                return "\n".join([p.extract_text() or "" for p in reader.pages])
            except Exception:
                # fall back to PyPDF2 if available
                try:
                    import PyPDF2
                    if use_bytes:
                        reader = PyPDF2.PdfReader(use_bytes)
                        return "\n".join([p.extract_text() or "" for p in reader.pages])
                    else:
                        with open(source, 'rb') as f:
                            reader = PyPDF2.PdfReader(f)
                            return "\n".join([p.extract_text() or "" for p in reader.pages])
                except Exception:
                    logger.exception('PDF extraction failed for %s', source or url)
                    return None

        # DOCX extraction
        if ext == '.docx':
            try:
                import docx
                if use_bytes:
                    doc = docx.Document(use_bytes)
                else:
                    doc = docx.Document(source)
                return "\n".join([p.text for p in doc.paragraphs])
            except Exception:
                logger.exception('DOCX extraction failed for %s', source or url)
                return None

        # PPTX extraction
        if ext == '.pptx':
            try:
                from pptx import Presentation
                if use_bytes:
                    prs = Presentation(use_bytes)
                else:
                    prs = Presentation(source)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            texts.append(shape.text)
                return "\n".join(texts)
            except Exception:
                logger.exception('PPTX extraction failed for %s', source or url)
                return None

    except Exception:
        logger.exception('Unexpected error extracting text from %s', source or url)
        return None

    return None

@api_view(['POST'])
def detect_submission(request):
    submission_id = request.data.get('submission_id')
    
    # 1. Get the submission
    try:
        submission = Submission.objects.get(id=submission_id)
    except Submission.DoesNotExist:
        return Response({"error": "Submission not found"}, status=404)
        
    text = submission.content or ''

    # If a file was uploaded, always try to extract its text and include it in analysis.
    if submission.file:
        extracted = _extract_text_from_file(submission.file)
        if extracted:
            if text.strip():
                text = f"{text}\n\n{extracted}"
            else:
                text = extracted
        elif not text.strip():
            report = {
                "is_ai_generated": False,
                "chunks": [],
                "note": "extraction_unavailable",
                "message": "No text found in submission and document extraction is unavailable."
            }
            DetectionResult.objects.update_or_create(
                submission=submission,
                defaults={
                    'is_ai_flagged': False,
                    'report_data': report
                }
            )
            return Response({"message": "Extraction failed", "report": report}, status=200)

    # 2. Run AI detection only; source lookup runs separately.
    report = run_analysis(text)
    
    # 3. Save to database
    result, created = DetectionResult.objects.update_or_create(
        submission=submission,
        defaults={
            'is_ai_flagged': report.get('is_ai_generated', False),
            'report_data': report
        }
    )
    
    return Response({
        "message": "Analysis saved", 
        "report": report
    })

@api_view(['POST'])
def verify_submission_sources(request):
    submission_id = request.data.get('submission_id')

    try:
        submission = Submission.objects.get(id=submission_id)
    except Submission.DoesNotExist:
        return Response({"error": "Submission not found"}, status=404)

    report = None
    existing = DetectionResult.objects.filter(submission=submission).order_by('-created_at').first()
    if existing and existing.report_data:
        report = existing.report_data
    else:
        text = submission.content or ''
        if submission.file:
            extracted = _extract_text_from_file(submission.file)
            if extracted:
                text = f"{text}\n\n{extracted}" if text.strip() else extracted
        report = run_analysis(text)

    def _verify_in_background(report_data):
        try:
            verified = run_source_verification(report_data)
            DetectionResult.objects.update_or_create(
                submission=submission,
                defaults={
                    'is_ai_flagged': verified.get('is_ai_generated', False),
                    'report_data': verified
                }
            )
        except Exception:
            logger.exception('Background source verification failed for submission %s', submission_id)

    Thread(target=_verify_in_background, args=(report,), daemon=True).start()

    return Response({
        "message": "Source verification started",
        "report": report
    })


@api_view(['POST'])
def run_batch_plagiarism_analysis(request):
    submission_ids = request.data.get('submission_ids') or []
    if not submission_ids:
        return Response({'error': 'At least one submission id is required.'}, status=400)

    submissions = []
    for submission_id in submission_ids:
        try:
            submission = Submission.objects.get(id=submission_id)
        except Submission.DoesNotExist:
            continue

        text = submission.content or ''
        if submission.file:
            extracted = extract_text_from_submission_file(submission.file)
            if extracted:
                text = f"{text}\n\n{extracted}" if text.strip() else extracted

        submissions.append({
            'id': submission.id,
            'title': submission.assignment.title if submission.assignment_id else 'Submission',
            'student_name': submission.student.username if getattr(submission.student, 'username', None) else '',
            'file_name': os.path.basename(submission.file.name) if submission.file else '',
            'file_url': submission.file.url if submission.file else '',
            'text': text,
        })

    if not submissions:
        return Response({'error': 'No valid submissions were found.'}, status=404)

    return Response(build_similarity_report(submissions))
