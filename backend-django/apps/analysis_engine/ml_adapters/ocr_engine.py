import io
import os
import re
import logging
import tempfile
from io import BytesIO
from typing import Optional, List

import requests
import shutil
import subprocess

try:
    from PIL import Image, ImageFilter, ImageOps
except ImportError:
    Image = None
    ImageFilter = None
    ImageOps = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    import easyocr
except ImportError:
    easyocr = None

try:
    import numpy as np
except ImportError:
    np = None

try:
    from transformers import TrOCRProcessor, VisionEncoderDecoderModel
except ImportError:
    TrOCRProcessor = None
    VisionEncoderDecoderModel = None

try:
    from pdf2image import convert_from_bytes, convert_from_path
except ImportError:
    convert_from_bytes = None
    convert_from_path = None

logger = logging.getLogger(__name__)

TROCR_MODEL_NAME = 'microsoft/trocr-base-handwritten'
_trocr_processor = None
_trocr_model = None
_reader = None

OCR_IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif'}
TEXT_FILE_EXTENSIONS = {'.txt'}


def _load_trocr_model() -> bool:
    global _trocr_processor, _trocr_model
    if _trocr_processor is not None and _trocr_model is not None:
        return True

    if TrOCRProcessor is None or VisionEncoderDecoderModel is None:
        logger.warning('TrOCR dependencies are not installed.')
        return False

    for use_fast in (True, False):
        try:
            _trocr_processor = TrOCRProcessor.from_pretrained(TROCR_MODEL_NAME, use_fast=use_fast)
            _trocr_model = VisionEncoderDecoderModel.from_pretrained(TROCR_MODEL_NAME)
            return True
        except Exception as exc:
            logger.exception('Failed to load TrOCR model with use_fast=%s: %s', use_fast, exc)
            _trocr_processor = None
            _trocr_model = None

    return False


def _load_easyocr_reader() -> bool:
    global _reader
    if _reader is not None:
        return True

    if easyocr is None or np is None:
        logger.warning('EasyOCR or numpy is not installed.')
        return False

    try:
        _reader = easyocr.Reader(['en'], gpu=False)
        return True
    except Exception as exc:
        logger.exception('Failed to create easyocr reader: %s', exc)
        _reader = None
        return False


def _load_image_from_bytes(data: BytesIO):
    if Image is None:
        return None
    try:
        data.seek(0)
        return Image.open(data).convert('RGB')
    except Exception:
        return None


def _clean_image(image, handwritten=False):
    if Image is None or image is None:
        return image

    try:
        gray = image.convert('L')
        gray = ImageOps.autocontrast(gray)
        if handwritten and ImageFilter is not None:
            gray = gray.filter(ImageFilter.MedianFilter(size=3))
            gray = gray.point(lambda p: 255 if p > 145 else 0)
        return gray
    except Exception:
        return image


def _aggressive_preprocess(image):
    """Aggressively preprocess a PIL Image to improve handwritten OCR recall.

    Steps:
    - Convert to grayscale
    - Contrast stretching (2-98 percentile)
    - Global thresholding using a conservative cutoff
    - Upscale small images
    - Unsharp mask to emphasize strokes
    """
    if Image is None or image is None or np is None:
        return image

    try:
        gray = image.convert('L')
        arr = np.array(gray).astype('float32')

        # Contrast stretch
        p2, p98 = np.percentile(arr, (2, 98))
        if p98 - p2 > 1:
            arr = np.clip((arr - p2) * 255.0 / (p98 - p2), 0, 255)

        # Slight blur to reduce noise then sharpen later
        proc = Image.fromarray(arr.astype('uint8'))
        if ImageFilter is not None:
            proc = proc.filter(ImageFilter.MedianFilter(size=3))

        # Global threshold (use mean * 0.9 to be slightly permissive)
        a = np.array(proc).astype('uint8')
        thresh = max(30, int(a.mean() * 0.9))
        binar = (a > thresh).astype('uint8') * 255
        proc = Image.fromarray(binar.astype('uint8'))

        # Upscale to help OCR
        w, h = proc.size
        if w < 1600 or h < 1600:
            new_w = min(2200, max(1600, w * 2))
            new_h = min(2200, max(1600, h * 2))
            proc = proc.resize((new_w, new_h), Image.BILINEAR)

        # Unsharp mask to emphasize pen strokes
        try:
            proc = proc.convert('L')
            proc = proc.filter(ImageFilter.UnsharpMask(radius=2, percent=150, threshold=3))
        except Exception:
            pass

        return proc.convert('RGB')
    except Exception:
        return image


def _generate_ocr_variants(image, handwritten=False):
    if Image is None or image is None:
        return []

    variants = []
    try:
        base = image.convert('RGB')
        variants.append(base)
        cleaned = _clean_image(base, handwritten=handwritten)
        if cleaned is not None:
            variants.append(cleaned.convert('RGB'))
            try:
                inverted = ImageOps.invert(cleaned)
                variants.append(inverted.convert('RGB'))
            except Exception:
                pass
        if base.mode != 'RGB':
            variants.append(base.convert('RGB'))
    except Exception:
        variants.append(image)

    # Remove duplicates by size/mode and preserve order
    unique = []
    keys = set()
    for img in variants:
        key = (img.size, img.mode)
        if key not in keys:
            keys.add(key)
            unique.append(img)
    return unique


def _run_tesseract(image, handwritten=False):
    if pytesseract is None or image is None:
        return None
    try:
        _find_and_configure_tesseract()
    except Exception:
        pass

    def attempt_tesseract(target_image, config):
        try:
            return pytesseract.image_to_string(target_image, lang='eng', config=config) or ''
        except Exception:
            logger.exception('Tesseract OCR attempt failed with config: %s', config)
            return ''

    best_text = ''
    # When handwriting is expected, run aggressive preprocessing first
    variants = []
    if handwritten:
        try:
            agg = _aggressive_preprocess(image)
            variants.extend(_generate_ocr_variants(agg, handwritten=True))
        except Exception:
            pass
    variants.extend(_generate_ocr_variants(image, handwritten=handwritten))

    for img in variants:
        for config in ['--oem 3 --psm 3', '--oem 3 --psm 6', '--oem 1 --psm 3', '--oem 1 --psm 6', '--oem 1 --psm 11', '--oem 1 --psm 4']:
            candidate = attempt_tesseract(img, config).strip()
            if candidate and len(candidate) > len(best_text):
                best_text = candidate
                if len(best_text) >= 100:
                    return best_text

    return best_text or None


def _run_easyocr(image):
    if image is None or easyocr is None or np is None:
        return None
    if not _load_easyocr_reader():
        return None

    try:
        # Aggressive preprocessing for handwritten scans
        proc_img = _aggressive_preprocess(image)
        try:
            proc_img = proc_img.convert('RGB')
        except Exception:
            proc_img = image.convert('RGB')

        # Upscale small scans to help EasyOCR
        try:
            w, h = proc_img.size
            if w < 1600 and h < 1600:
                proc_img = proc_img.resize((min(2000, w * 2), min(2000, h * 2)))
        except Exception:
            pass

        texts = _reader.readtext(np.array(proc_img), detail=0)
        if texts:
            return '\n'.join([t.strip() for t in texts if t and t.strip()]).strip() or None
    except Exception:
        logger.exception('EasyOCR failed.')
    return None


def _find_and_configure_tesseract() -> bool:
    """Try to find a Tesseract binary on the system and configure pytesseract.

    Returns True if tesseract executable is found and configured, otherwise False.
    """
    if pytesseract is None:
        return False

    try:
        cmd = getattr(pytesseract.pytesseract, 'tesseract_cmd', None)
        if cmd:
            if shutil.which(cmd) or os.path.exists(cmd):
                return True
    except Exception:
        pass

    candidates = [
        shutil.which('tesseract'),
        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
    ]

    for c in candidates:
        if not c:
            continue
        try:
            if shutil.which(c) or os.path.exists(c):
                pytesseract.pytesseract.tesseract_cmd = c
                return True
        except Exception:
            continue

    return False


def _run_google_vision(image):
    """Use Google Cloud Vision as a fallback if available. Returns text or None."""
    try:
        from google.cloud import vision
    except Exception:
        return None

    try:
        client = vision.ImageAnnotatorClient()
        buffered = BytesIO()
        image.save(buffered, format='PNG')
        content = buffered.getvalue()
        gimg = vision.Image(content=content)
        response = client.document_text_detection(image=gimg)
        if response.error.message:
            logger.warning('Google Vision error: %s', response.error.message)
            return None
        annotation = response.full_text_annotation
        if annotation and annotation.text:
            return annotation.text.strip() or None
        if response.text_annotations:
            return response.text_annotations[0].description.strip() or None
    except Exception:
        logger.exception('Google Vision OCR failed.')
    return None


def _run_trocr(image):
    if not _load_trocr_model() or image is None:
        return None
    try:
        processed = _trocr_processor(images=image, return_tensors='pt')
        generated_ids = _trocr_model.generate(processed.pixel_values, max_length=512)
        return _trocr_processor.batch_decode(generated_ids, skip_special_tokens=True)[0]
    except Exception:
        logger.exception('TrOCR inference failed.')
        return None


def _choose_best_text(tesseract_text: Optional[str], trocr_text: Optional[str], easyocr_text: Optional[str] = None) -> Optional[str]:
    return _choose_best_text_full(tesseract_text, trocr_text, easyocr_text, None)


def _choose_best_text_full(tesseract_text: Optional[str], trocr_text: Optional[str], easyocr_text: Optional[str] = None, google_text: Optional[str] = None) -> Optional[str]:
    tesseract_text = (tesseract_text or '').strip()
    trocr_text = (trocr_text or '').strip()
    easyocr_text = (easyocr_text or '').strip()
    google_text = (google_text or '').strip()

    candidates = [t for t in [tesseract_text, trocr_text, easyocr_text, google_text] if t]
    if not candidates:
        return None
    return max(candidates, key=len)


def _extract_text_from_image(image):
    """Extract text from a PIL Image using available OCR engines.

    Strategy:
    - Run Tesseract first (fast). If Tesseract output is short/weak, prefer EasyOCR.
    - Always attempt TroCR and Google Vision (if available) as fallbacks.
    - Prefer the longest reasonable candidate, but give EasyOCR priority when
      Tesseract produced very little text.
    """
    if image is None:
        return None

    try:
        tesseract_text = _run_tesseract(image)
    except Exception:
        tesseract_text = None

    try:
        easyocr_text = _run_easyocr(image)
    except Exception:
        easyocr_text = None

    try:
        trocr_text = _run_trocr(image)
    except Exception:
        trocr_text = None

    try:
        google_text = _run_google_vision(image)
    except Exception:
        google_text = None

    t_len = len((tesseract_text or '').strip())
    e_len = len((easyocr_text or '').strip())

    # If Tesseract is weak but EasyOCR produced more, prefer EasyOCR.
    if t_len < 60 and e_len > t_len:
        return _choose_best_text_full(easyocr_text, trocr_text, None, google_text)

    return _choose_best_text_full(tesseract_text, trocr_text, easyocr_text, google_text)


def _get_file_bytes(file_field):
    if not file_field:
        return None, None, None

    storage = getattr(file_field, 'storage', None)
    name = getattr(file_field, 'name', None)
    path = getattr(file_field, 'path', None)
    url = getattr(file_field, 'url', None)

    use_bytes = None
    source = None
    if path and os.path.exists(path):
        source = path
    elif storage and name:
        try:
            if storage.exists(name):
                with storage.open(name, 'rb') as handle:
                    use_bytes = BytesIO(handle.read())
        except Exception:
            logger.exception('Failed to open file from storage: %s', name)

    if source is None and use_bytes is None:
        if hasattr(file_field, 'open'):
            try:
                file_field.open('rb')
                use_bytes = BytesIO(file_field.read())
            except Exception:
                logger.exception('Failed to read FileField contents for extraction: %s', name or url)
            finally:
                try:
                    file_field.close()
                except Exception:
                    pass

    if source is None and use_bytes is None and hasattr(file_field, 'read'):
        try:
            file_field.seek(0)
        except Exception:
            pass
        try:
            use_bytes = BytesIO(file_field.read())
        except Exception:
            logger.exception('Failed to read file-like object for extraction: %s', name or url)

    if source is None and use_bytes is None:
        file_obj = getattr(file_field, '_file', None)
        if file_obj is not None:
            try:
                file_obj.seek(0)
                use_bytes = BytesIO(file_obj.read())
            except Exception:
                logger.exception('Failed to read file-like object for extraction: %s', name or url)

    if source is None and use_bytes is None and url and url.startswith(('http://', 'https://')):
        try:
            response = requests.get(url, timeout=10)
            if response.status_code == 200:
                use_bytes = BytesIO(response.content)
        except Exception:
            logger.exception('Failed to download file from URL: %s', url)

    ext = None
    if source:
        ext = os.path.splitext(source)[1].lower()
    elif name:
        ext = os.path.splitext(name.split('?')[0])[1].lower()
    elif url:
        ext = os.path.splitext(url.split('?')[0])[1].lower()

    return source, use_bytes, ext


def _load_pdf_images(source, use_bytes):
    images = []
    if Image is None:
        return images

    if convert_from_bytes is not None and convert_from_path is not None:
        try:
            if use_bytes is not None:
                use_bytes.seek(0)
                images = convert_from_bytes(use_bytes.getvalue(), dpi=300, fmt='png')
            elif source is not None:
                images = convert_from_path(source, dpi=300, fmt='png')
            if images:
                return images
        except Exception:
            logger.exception('Failed to convert PDF to images for OCR.')

    try:
        import fitz
        if use_bytes is not None:
            use_bytes.seek(0)
            doc = fitz.open(stream=use_bytes.getvalue(), filetype='pdf')
        elif source is not None:
            doc = fitz.open(source)
        else:
            return images

        for page in doc:
            pix = page.get_pixmap(dpi=300)
            image_data = BytesIO(pix.tobytes('png'))
            img = Image.open(image_data).convert('RGB')
            images.append(img)
        return images
    except Exception:
        logger.exception('PyMuPDF PDF rendering failed.')

    try:
        pdftoppm = shutil.which('pdftoppm')
        if not pdftoppm:
            return images

        with tempfile.TemporaryDirectory() as tmpdir:
            if use_bytes is not None:
                temp_pdf = os.path.join(tmpdir, 'temp.pdf')
                use_bytes.seek(0)
                with open(temp_pdf, 'wb') as handle:
                    handle.write(use_bytes.getvalue())
            elif source is not None:
                temp_pdf = source
            else:
                return images

            cmd = [pdftoppm, '-png', '-r', '300', temp_pdf, os.path.join(tmpdir, 'page')]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

            for file_name in sorted(os.listdir(tmpdir)):
                if file_name.startswith('page') and file_name.lower().endswith('.png'):
                    img_path = os.path.join(tmpdir, file_name)
                    try:
                        img = Image.open(img_path).convert('RGB')
                        images.append(img)
                    except Exception:
                        logger.exception('Failed to open pdftoppm image %s', img_path)
        return images
    except Exception:
        logger.exception('pdftoppm conversion failed.')

    return images


def _extract_text_from_pdf(source, use_bytes):
    images = _load_pdf_images(source, use_bytes)
    pages = []
    for img in images:
        page_text = _extract_text_from_image(img)
        if page_text:
            page_text = _normalize_extracted_text(page_text)
            if page_text:
                pages.append(page_text)
    return '\n\n'.join(pages).strip() if pages else None


def _normalize_extracted_text(text: str) -> str:
    # Normalize line breaks and reduce OCR fragmentation.
    normalized = text.replace('\r\n', '\n').replace('\r', '\n')
    normalized = re.sub(r'[ \t]+', ' ', normalized)
    normalized = re.sub(r'(?<!\n)\n(?!\n)', ' ', normalized)
    normalized = re.sub(r'\n{2,}', '\n\n', normalized)
    normalized = re.sub(r' {2,}', ' ', normalized)
    return normalized.strip()


def _extract_text_from_image_bytes(data: BytesIO):
    image = _load_image_from_bytes(data)
    if image is None:
        return None
    text = _extract_text_from_image(image)
    if text:
        return _normalize_extracted_text(text)
    return None


def _extract_text_from_text_file(source, use_bytes, ext):
    if ext == '.docx':
        try:
            import docx
            document = docx.Document(use_bytes) if use_bytes else docx.Document(source)
            return '\n'.join([p.text for p in document.paragraphs if p.text])
        except Exception:
            logger.exception('DOCX extraction failed.')
            return None

    if ext == '.pptx':
        try:
            from pptx import Presentation
            presentation = Presentation(use_bytes) if use_bytes else Presentation(source)
            texts = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        texts.append(shape.text)
            return '\n'.join(texts)
        except Exception:
            logger.exception('PPTX extraction failed.')
            return None

    if ext in {'.html', '.htm'}:
        try:
            from html.parser import HTMLParser

            class HTMLTextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text = []
                    self.skip_content = False

                def handle_starttag(self, tag, attrs):
                    if tag in ['script', 'style', 'head', 'meta']:
                        self.skip_content = True

                def handle_endtag(self, tag):
                    if tag in ['script', 'style', 'head', 'meta']:
                        self.skip_content = False

                def handle_data(self, data):
                    if not self.skip_content:
                        text = data.strip()
                        if text:
                            self.text.append(text)

                def get_text(self):
                    return '\n'.join(self.text)

            if use_bytes:
                html_content = use_bytes.read().decode('utf-8', errors='ignore')
            else:
                with open(source, 'r', encoding='utf-8', errors='ignore') as handle:
                    html_content = handle.read()

            parser = HTMLTextExtractor()
            parser.feed(html_content)
            return parser.get_text()
        except Exception:
            logger.exception('HTML extraction failed.')
            return None

    if ext == '.txt':
        try:
            if use_bytes:
                use_bytes.seek(0)
                return use_bytes.read().decode('utf-8', errors='ignore')
            with open(source, 'r', encoding='utf-8', errors='ignore') as handle:
                return handle.read()
        except Exception:
            logger.exception('Text extraction failed.')
            return None

    return None


def extract_text_from_file(file_field, return_extraction_type=False):
    source, use_bytes, ext = _get_file_bytes(file_field)
    if not ext:
        return None if not return_extraction_type else (None, False)

    ocr_extraction = False
    extracted_text = None

    if ext == '.pdf':
        if source or use_bytes:
            text = _extract_text_from_pdf(source, use_bytes)
            if text:
                extracted_text = text
                ocr_extraction = True
            else:
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(use_bytes) if use_bytes else PdfReader(source)
                    extracted = '\n'.join([p.extract_text() or '' for p in reader.pages]).strip()
                    if extracted:
                        extracted_text = extracted
                        ocr_extraction = False
                except Exception:
                    logger.exception('PDF text extraction failed with pypdf.')
                if extracted_text is None:
                    try:
                        import PyPDF2
                        reader = PyPDF2.PdfReader(use_bytes) if use_bytes else PyPDF2.PdfReader(source)
                        extracted = '\n'.join([p.extract_text() or '' for p in reader.pages]).strip()
                        if extracted:
                            extracted_text = extracted
                            ocr_extraction = False
                    except Exception:
                        logger.exception('PDF text extraction failed with PyPDF2.')
        if extracted_text:
            return extracted_text if not return_extraction_type else (extracted_text, ocr_extraction)
        return None if not return_extraction_type else (None, False)

    if ext in OCR_IMAGE_EXTENSIONS:
        text = _extract_text_from_image_bytes(use_bytes or BytesIO(open(source, 'rb').read()))
        if text:
            return (text, True) if return_extraction_type else text

    text = _extract_text_from_text_file(source, use_bytes, ext)
    if text:
        return (text, False) if return_extraction_type else text

    return None if not return_extraction_type else (None, False)
