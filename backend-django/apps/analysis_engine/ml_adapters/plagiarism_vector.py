import math
import os
import re
from collections import Counter
from io import BytesIO

import requests
from django.core.files.base import File

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
except Exception:  # pragma: no cover - optional dependency fallback
    TfidfVectorizer = None
    cosine_similarity = None

"""
Plagiarism Detection Module

This module compares submissions using three complementary similarity metrics:

1. **Jaccard Similarity (J)**: Set-based token overlap
   - Counts unique tokens (words) in both documents
   - Formula: |intersection| / |union|
   - Range: 0 to 1 (higher = more similar)
   - Best for: Detecting similar vocabulary/concepts
   
2. **TF-IDF Similarity (T)**: Term Frequency-Inverse Document Frequency
   - Uses sklearn's TfidfVectorizer with bigrams
   - Measures importance of terms in documents
   - Formula: Cosine similarity between TF-IDF vectors
   - Range: 0 to 1 (higher = more similar)
   - Best for: Detecting content similarity while reducing common words
   
3. **Semantic Similarity (S)**: Weighted token overlap
   - Considers frequency of shared tokens
   - Normalized by document term frequency vectors
   - Formula: sum(min_frequency) / sqrt(left_norm² × right_norm²)
   - Range: 0 to 1 (higher = more similar)
   - Best for: Finding substantial shared content

Chunking Strategy (same as AI detector):
- Split text into paragraphs (double newlines)
- Split paragraphs into sentences (., !, ?)
- Group sentences into chunks:
  * Paragraphs >= 8 sentences: split into 4-sentence chunks
  * Paragraphs < 3 sentences: buffer together
  * Paragraphs 3-7 sentences: create single chunk or buffer
"""

SENTENCE_SPLIT_PATTERN = re.compile(r'(?<=[.!?])\s+(?=[A-Z0-9])')


def _split_paragraphs(text):
    """Split text into paragraphs by double newlines."""
    normalized = text.replace('\r\n', '\n').replace('\r', '\n')
    paragraphs = [p.strip() for p in re.split(r'\n{2,}', normalized) if p.strip()]
    return paragraphs if paragraphs else [normalized.strip()]


def _split_sentences(paragraph):
    """Split paragraph into sentences."""
    sentences = [s.strip() for s in SENTENCE_SPLIT_PATTERN.split(paragraph) if s.strip()]
    return sentences if sentences else [paragraph.strip()]


# --- Chunking configuration -------------------------------------------------
# Fallback chunk sizing when real paragraphs are not recognizable.
SENTENCES_PER_CHUNK = 4          # ~4 sentences per generated chunk
WORDS_PER_CHUNK = 100            # ...or ~100 words, whichever comes first
# Overlap window used ONLY to give each generated chunk more context for
# scoring. The overlap is NEVER highlighted/displayed.
OVERLAP_SENTENCES = 2            # 1-2 sentences of context above/below
OVERLAP_WORDS = 40               # ~25-50 words of context above/below
# A recognizable paragraph larger than this many sentences is treated as an
# unstructured block and re-chunked with the sentence/word fallback.
LARGE_PARAGRAPH_SENTENCES = 8


def _fallback_sentence_groups(sentences):
    """Group a flat list of sentences into core windows of ~4 sentences /
    ~100 words. Returns a list of (start_idx, end_idx) tuples (end exclusive)."""
    groups = []
    n = len(sentences)
    i = 0
    while i < n:
        start = i
        word_count = 0
        count = 0
        while i < n and count < SENTENCES_PER_CHUNK and word_count < WORDS_PER_CHUNK:
            word_count += len(sentences[i].split())
            count += 1
            i += 1
        groups.append((start, i))
    return groups


def _context_prefix(sentences, start):
    """Build overlap context (before the core) limited by sentence/word caps."""
    words_used = 0
    picked = []
    for idx in range(start - 1, -1, -1):
        if len(picked) >= OVERLAP_SENTENCES or words_used >= OVERLAP_WORDS:
            break
        picked.append(sentences[idx])
        words_used += len(sentences[idx].split())
    return ' '.join(reversed(picked)).strip()


def _context_suffix(sentences, end):
    """Build overlap context (after the core) limited by sentence/word caps."""
    words_used = 0
    picked = []
    for idx in range(end, len(sentences)):
        if len(picked) >= OVERLAP_SENTENCES or words_used >= OVERLAP_WORDS:
            break
        picked.append(sentences[idx])
        words_used += len(sentences[idx].split())
    return ' '.join(picked).strip()


def _build_chunks_from_paragraphs(paragraphs):
    """Paragraph-first chunking.

    Rules implemented (per project requirement):
    - Each recognizable paragraph is treated as ONE chunk.
    - When paragraphs/points are not recognizable (single-paragraph document)
      OR a paragraph is an unusually large block, split it into generated
      chunks of ~4 sentences / ~100 words.
    - For those generated chunks, an overlap window of 1-2 sentences / 25-50
      words above and below is attached to the ``compare_text`` so the chunk
      keeps its surrounding meaning during scoring. The overlap is deliberately
      excluded from ``core_text`` so highlighting/display only covers the real
      chunk content.

    Returns a list of dicts: ``{'core_text', 'compare_text'}``.
    """
    chunks = []
    single_paragraph_doc = len(paragraphs) == 1

    for paragraph in paragraphs:
        sentences = _split_sentences(paragraph)
        if not sentences:
            continue

        recognizable_paragraph = not single_paragraph_doc and len(sentences) <= LARGE_PARAGRAPH_SENTENCES

        if recognizable_paragraph:
            # Whole paragraph = one chunk, no overlap needed.
            text = paragraph.strip()
            chunks.append({'core_text': text, 'compare_text': text})
            continue

        # Unstructured / large block: generate 4-sentence / 100-word chunks
        # with an overlap window used only for scoring context.
        for start, end in _fallback_sentence_groups(sentences):
            core_text = ' '.join(sentences[start:end]).strip()
            if not core_text:
                continue
            prefix = _context_prefix(sentences, start)
            suffix = _context_suffix(sentences, end)
            compare_text = ' '.join(part for part in (prefix, core_text, suffix) if part).strip()
            chunks.append({'core_text': core_text, 'compare_text': compare_text})

    return chunks


def extract_text_from_submission_file(file_field):
    """Extract text from PDF, DOCX, or PPTX files from storage or a remote URL."""
    if not file_field:
        return None

    storage = getattr(file_field, 'storage', None)
    name = getattr(file_field, 'name', None)
    path = getattr(file_field, 'path', None)
    url = getattr(file_field, 'url', None)

    use_bytes = None
    source = None

    if path and os.path.exists(path):
        source = path
    elif storage and name:
        try:
            if storage.exists(name):
                with storage.open(name, 'rb') as handle:
                    use_bytes = BytesIO(handle.read())
        except Exception:
            use_bytes = None

    if source is None and use_bytes is None and url and (url.startswith('http://') or url.startswith('https://')):
        try:
            response = requests.get(url, timeout=10)
            if response.status_code == 200:
                use_bytes = BytesIO(response.content)
        except Exception:
            use_bytes = None

    ext = None
    if source:
        ext = os.path.splitext(source)[1].lower()
    elif name:
        ext = os.path.splitext(name.split('?')[0])[1].lower()
    elif url:
        ext = os.path.splitext(url.split('?')[0])[1].lower()

    if ext == '.pdf':
        try:
            from pypdf import PdfReader
            reader = PdfReader(use_bytes) if use_bytes else PdfReader(source)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(use_bytes) if use_bytes else PyPDF2.PdfReader(source)
                return "\n".join(page.extract_text() or "" for page in reader.pages)
            except Exception:
                return None

    if ext == '.docx':
        try:
            import docx
            document = docx.Document(use_bytes) if use_bytes else docx.Document(source)
            return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
        except Exception:
            return None

    if ext == '.pptx':
        try:
            from pptx import Presentation
            presentation = Presentation(use_bytes) if use_bytes else Presentation(source)
            texts = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception:
            return None

    if ext == '.html' or ext == '.htm':
        try:
            from html.parser import HTMLParser
            
            class HTMLTextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text = []
                    self.skip_content = False

                def handle_starttag(self, tag, attrs):
                    if tag in ['script', 'style', 'head', 'meta']:
                        self.skip_content = True

                def handle_endtag(self, tag):
                    if tag in ['script', 'style', 'head', 'meta']:
                        self.skip_content = False

                def handle_data(self, data):
                    if not self.skip_content:
                        text = data.strip()
                        if text:
                            self.text.append(text)

                def get_text(self):
                    return "\n".join(self.text)

            if use_bytes:
                html_content = use_bytes.read().decode('utf-8', errors='ignore')
            else:
                with open(source, 'r', encoding='utf-8', errors='ignore') as f:
                    html_content = f.read()
            
            parser = HTMLTextExtractor()
            parser.feed(html_content)
            return parser.get_text()
        except Exception:
            return None

    return None


def _tokenize(text):
    return [token for token in re.findall(r"[a-zA-Z0-9']+", text.lower()) if token]


def _normalize_text(text):
    return " ".join(_tokenize(text))


def _jaccard_similarity(left_text, right_text):
    left_tokens = set(_tokenize(left_text))
    right_tokens = set(_tokenize(right_text))
    if not left_tokens and not right_tokens:
        return 0.0
    union = left_tokens | right_tokens
    if not union:
        return 0.0
    return round(len(left_tokens & right_tokens) / len(union), 3)


def _tfidf_similarity(left_text, right_text):
    if TfidfVectorizer is None or cosine_similarity is None:
        return _jaccard_similarity(left_text, right_text)

    vectorizer = TfidfVectorizer(ngram_range=(1, 2), stop_words='english')
    matrix = vectorizer.fit_transform([left_text, right_text])
    score = cosine_similarity(matrix[0:1], matrix[1:2])[0][0]
    return round(float(score), 3)


def _semantic_similarity(left_text, right_text):
    left_tokens = Counter(_tokenize(left_text))
    right_tokens = Counter(_tokenize(right_text))
    if not left_tokens or not right_tokens:
        return 0.0

    common_tokens = set(left_tokens.keys()) & set(right_tokens.keys())
    if not common_tokens:
        return 0.0

    numerator = sum(min(left_tokens[token], right_tokens[token]) for token in common_tokens)
    left_norm = math.sqrt(sum(value * value for value in left_tokens.values()))
    right_norm = math.sqrt(sum(value * value for value in right_tokens.values()))
    denominator = left_norm * right_norm
    if denominator == 0:
        return 0.0
    return round(numerator / denominator, 3)


def _comparison_verdict(score):
    if score >= 0.65:
        return 'High similarity'
    if score >= 0.4:
        return 'Moderate similarity'
    if score >= 0.2:
        return 'Low similarity'
    return 'Clear'


def _chunk_core(chunk):
    """Return the display/highlight text for a chunk (no overlap window)."""
    if isinstance(chunk, dict):
        return chunk.get('core_text', '')
    return chunk


def _chunk_compare(chunk):
    """Return the scoring text for a chunk (may include overlap context)."""
    if isinstance(chunk, dict):
        return chunk.get('compare_text') or chunk.get('core_text', '')
    return chunk


def _compare_chunks(left_chunks, right_chunks):
    """Compare chunks between two documents and return detailed chunk-level analysis.

    Scoring is performed on each chunk's ``compare_text`` (which may include the
    overlap window for context), while the stored preview text uses the
    ``core_text`` so downstream highlighting only covers the real chunk content.

    Returns:
        - chunk_matches: List of matching chunks with their scores and color IDs
    """
    candidates = []

    for left_idx, left_chunk in enumerate(left_chunks):
        left_compare = _chunk_compare(left_chunk)
        left_core = _chunk_core(left_chunk)
        for right_idx, right_chunk in enumerate(right_chunks):
            right_compare = _chunk_compare(right_chunk)
            right_core = _chunk_core(right_chunk)
            jaccard = _jaccard_similarity(left_compare, right_compare)
            tfidf = _tfidf_similarity(left_compare, right_compare)
            semantic = _semantic_similarity(left_compare, right_compare)
            is_match = jaccard >= 0.15 or tfidf >= 0.15 or semantic >= 0.12
            if not is_match:
                continue

            score = round((jaccard + tfidf + semantic) / 3, 3)
            candidates.append({
                'left_chunk_idx': left_idx,
                'right_chunk_idx': right_idx,
                'left_text': left_core[:200],
                'right_text': right_core[:200],
                'scores': {
                    'jaccard': jaccard,
                    'tfidf': tfidf,
                    'semantic': semantic,
                },
                'score': score,
            })

    # Choose the strongest unique left/right matches to avoid duplicate chunk reuse.
    candidates.sort(key=lambda item: item['score'], reverse=True)
    assigned_left = set()
    assigned_right = set()
    chunk_matches = []

    for candidate in candidates:
        left_idx = candidate['left_chunk_idx']
        right_idx = candidate['right_chunk_idx']
        if left_idx in assigned_left or right_idx in assigned_right:
            continue

        candidate['is_match'] = True
        candidate['color_id'] = len(chunk_matches)
        chunk_matches.append(candidate)
        assigned_left.add(left_idx)
        assigned_right.add(right_idx)

    return chunk_matches


def _find_chunk_positions(text, chunks):
    """Find character positions of chunks in the original text."""
    positions = []
    search_start = 0
    
    for chunk in chunks:
        # Find the chunk in the text (with some tolerance for whitespace)
        chunk_words = chunk.split()
        if not chunk_words:
            positions.append({'start': 0, 'end': 0})
            continue
        
        # Find position of first word in chunk
        first_word = chunk_words[0]
        pos = text.find(first_word, search_start)
        
        if pos >= 0:
            # Find end by searching for last word
            last_word = chunk_words[-1]
            end_pos = text.find(last_word, pos)
            if end_pos >= 0:
                end_pos += len(last_word)
                positions.append({'start': pos, 'end': end_pos})
                search_start = end_pos
            else:
                positions.append({'start': pos, 'end': pos + len(first_word)})
                search_start = pos + len(first_word)
        else:
            positions.append({'start': 0, 'end': 0})
    
    return positions


def _build_chunk_highlights(left_text, left_chunks, right_text, right_chunks, chunk_matches):
    """Build highlight ranges for matching chunks.

    Positions are computed from each chunk's CORE text only, so the highlight in
    the original document covers exactly the chunk content and never the overlap
    context window used for scoring.
    """
    left_positions = _find_chunk_positions(left_text, [_chunk_core(c) for c in left_chunks])
    right_positions = _find_chunk_positions(right_text, [_chunk_core(c) for c in right_chunks])
    
    highlights = []
    
    for match in chunk_matches:
        left_idx = match['left_chunk_idx']
        right_idx = match['right_chunk_idx']
        color_id = match.get('color_id', -1)
        
        left_pos = left_positions[left_idx] if left_idx < len(left_positions) else {'start': 0, 'end': 0}
        right_pos = right_positions[right_idx] if right_idx < len(right_positions) else {'start': 0, 'end': 0}
        
        highlights.append({
            'left': left_pos if left_pos['start'] < left_pos['end'] else None,
            'right': right_pos if right_pos['start'] < right_pos['end'] else None,
            'color_id': color_id,
            'text': match['left_text'],
            'jaccard': match['scores']['jaccard'],
            'tfidf': match['scores']['tfidf'],
            'semantic': match['scores']['semantic'],
        })
    
    return highlights


def _sentence_spans(text):
    """Split text into sentences, returning each sentence with its char start/end
    positions in the ORIGINAL text so highlighting preserves the exact arrangement."""
    spans = []
    if not text:
        return spans
    # Match sentences ending with ., !, ? or the final trailing segment
    for match in re.finditer(r'[^.!?\n]+[.!?]?', text):
        segment = match.group()
        if segment.strip():
            spans.append({
                'text': segment.strip(),
                'start': match.start() + (len(segment) - len(segment.lstrip())),
                'end': match.end() - (len(segment) - len(segment.rstrip())),
            })
    return spans


def _build_sentence_highlights(left_text, right_text):
    """Match sentences between the two documents and assign each matched pair a
    distinct color id. Reordered sentences still match because comparison is
    all-pairs and position-independent. Returns highlight ranges preserving the
    original text arrangement, each carrying its own color + per-pair scores."""
    left_spans = _sentence_spans(left_text)
    right_spans = _sentence_spans(right_text)

    highlights = []
    color_id = 0
    used_right = set()

    for left_span in left_spans:
        left_tokens = set(_tokenize(left_span['text']))
        if not left_tokens:
            continue

        best = None
        best_score = 0.0
        for r_idx, right_span in enumerate(right_spans):
            if r_idx in used_right:
                continue
            right_tokens = set(_tokenize(right_span['text']))
            if not right_tokens:
                continue
            overlap = len(left_tokens & right_tokens)
            if overlap == 0:
                continue
            union = len(left_tokens | right_tokens)
            jaccard = overlap / union if union else 0.0
            if jaccard > best_score:
                best_score = jaccard
                best = (r_idx, right_span, right_tokens)

        # Threshold: sentences sharing enough tokens are considered similar chunks
        if best and best_score >= 0.3:
            r_idx, right_span, right_tokens = best
            used_right.add(r_idx)

            semantic = _semantic_similarity(left_span['text'], right_span['text'])
            tfidf = _tfidf_similarity(left_span['text'], right_span['text'])

            highlights.append({
                'left': {'start': left_span['start'], 'end': left_span['end']},
                'right': {'start': right_span['start'], 'end': right_span['end']},
                'color_id': color_id,
                'text': left_span['text'][:200],
                'jaccard': round(best_score, 3),
                'tfidf': round(tfidf, 3),
                'semantic': round(semantic, 3),
            })
            color_id += 1

    return highlights


def _extract_sentences(text):
    if not text:
        return []
    sentences = [sentence.strip() for sentence in re.split(r'(?<=[.!?])\s+', text) if sentence.strip()]
    return sentences[:8]



def _find_shared_snippets(left_text, right_text):
    left_sentences = _extract_sentences(left_text)
    right_sentences = _extract_sentences(right_text)
    shared = []

    for left_sentence in left_sentences:
        left_tokens = set(_tokenize(left_sentence))
        if not left_tokens:
            continue
        for right_sentence in right_sentences:
            right_tokens = set(_tokenize(right_sentence))
            if not right_tokens:
                continue
            overlap = len(left_tokens & right_tokens)
            if overlap == 0:
                continue
            score = overlap / len(left_tokens | right_tokens)
            if score >= 0.25:
                shared.append({
                    'left_snippet': left_sentence,
                    'right_snippet': right_sentence,
                    'overlap_score': round(score, 3),
                })

    return shared[:3]


def _build_highlights(left_text, right_text, shared_snippets):
    highlights = []
    for snippet in shared_snippets:
        left_range = None
        right_range = None
        left_phrase = snippet['left_snippet']
        right_phrase = snippet['right_snippet']
        if left_phrase:
            left_start = left_text.lower().find(left_phrase.lower())
            if left_start >= 0:
                left_range = {'start': left_start, 'end': left_start + len(left_phrase)}
        if right_phrase:
            right_start = right_text.lower().find(right_phrase.lower())
            if right_start >= 0:
                right_range = {'start': right_start, 'end': right_start + len(right_phrase)}
        if left_range or right_range:
            highlights.append({
                'left': left_range,
                'right': right_range,
                'text': left_phrase,
                'overlap_score': snippet['overlap_score'],
            })
    return highlights


def build_similarity_report(submissions):
    """Build a pairwise plagiarism report using chunk-based comparison.
    
    Chunks matching chunks on the other file are highlighted with the same color.
    Per-chunk similarity scores show how similar the paired chunks are.
    """
    matrix = []
    for left_index, left_submission in enumerate(submissions):
        row = []
        for right_index, right_submission in enumerate(submissions):
            if left_index == right_index:
                row.append({
                    'submission_id': left_submission['id'],
                    'submission_title': left_submission.get('title') or left_submission.get('student_name') or 'Submission',
                    'scores': {'jaccard': 1.0, 'tfidf': 1.0, 'semantic': 1.0},
                    'flagged': False,
                    'chunk_matches': [],
                    'highlights': [],
                    'is_diagonal': True,
                })
                continue

            left_text = left_submission.get('text') or ''
            right_text = right_submission.get('text') or ''
            
            # Build chunks using the same logic as AI detector
            left_paragraphs = _split_paragraphs(left_text)
            right_paragraphs = _split_paragraphs(right_text)
            left_chunks = _build_chunks_from_paragraphs(left_paragraphs)
            right_chunks = _build_chunks_from_paragraphs(right_paragraphs)
            
            # Compare chunks (kept for scoring / flagging)
            chunk_matches = _compare_chunks(left_chunks, right_chunks)

            # Build chunk-based highlights first, which preserve the same color
            # for matching chunks on both sides. Fall back to sentence-level
            # matches only when chunk-level matches are absent.
            highlights = _build_chunk_highlights(
                left_text,
                left_chunks,
                right_text,
                right_chunks,
                chunk_matches,
            ) if chunk_matches else _build_sentence_highlights(left_text, right_text)

            # Calculate overall document similarity — always use full-document metrics
            # so reordered/paraphrased content is still detected accurately
            doc_jaccard = _jaccard_similarity(left_text, right_text)
            doc_tfidf = _tfidf_similarity(left_text, right_text)
            doc_semantic = _semantic_similarity(left_text, right_text)

            # If chunk matches exist, take the max of doc-level and best-chunk scores
            if chunk_matches:
                best_jaccard = max(m['scores']['jaccard'] for m in chunk_matches)
                best_tfidf = max(m['scores']['tfidf'] for m in chunk_matches)
                best_semantic = max(m['scores']['semantic'] for m in chunk_matches)
                jaccard = round(max(doc_jaccard, best_jaccard), 3)
                tfidf = round(max(doc_tfidf, best_tfidf), 3)
                semantic = round(max(doc_semantic, best_semantic), 3)
            else:
                jaccard = round(doc_jaccard, 3)
                tfidf = round(doc_tfidf, 3)
                semantic = round(doc_semantic, 3)
            
            overall_score = round((jaccard + tfidf + semantic) / 3, 3)
            verdict = _comparison_verdict(overall_score)
            flagged = any([
                overall_score >= 0.2,
                bool(chunk_matches),
            ])

            row.append({
                'submission_id': right_submission['id'],
                'submission_title': right_submission.get('title') or right_submission.get('student_name') or 'Submission',
                'scores': {'jaccard': jaccard, 'tfidf': tfidf, 'semantic': semantic},
                'overall_score': overall_score,
                'verdict': verdict,
                'flagged': flagged,
                'chunk_matches': chunk_matches,
                'highlights': highlights,
                'is_diagonal': False,
            })
        matrix.append(row)

    return {
        'matrix': matrix,
        'summary': {
            'selected_count': len(submissions),
            'flagged_pairs': sum(1 for row in matrix for cell in row if cell.get('flagged') and not cell.get('is_diagonal')),
        },
        'submitted_files': [
            {
                'id': submission['id'],
                'title': submission.get('title') or submission.get('student_name') or 'Submission',
                'student_name': submission.get('student_name') or '',
                'file_name': submission.get('file_name') or '',
                'file_url': submission.get('file_url') or '',
                'text': submission.get('text') or '',
            }
            for submission in submissions
        ],
    }
